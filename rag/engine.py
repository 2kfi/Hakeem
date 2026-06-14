import asyncio
import logging
import os
from pathlib import Path
from typing import Any, Optional

import numpy as np

from core.config import HakeemRAGSettings
from rag.download import download_onnx_model
from rag.hakeem_citation import HakeemCitationFormatter
from rag.hakeem_corrective_rag import HakeemCorrectiveRAG
from rag.hakeem_hybrid_retriever import HakeemHybridRetriever
from rag.hakeem_knowledge_graph import HakeemKnowledgeGraph
from rag.hakeem_parent_retriever import HakeemParentRetriever
from rag.hakeem_qdrant_store import HakeemQdrantStore
from rag.hakeem_query_decomposer import HakeemQueryDecomposer
from rag.hakeem_reranker import HakeemReranker
from rag.hakeem_semantic_router import HakeemSemanticRouter
from rag.onnx_embedding import OnnxEmbedding, _load_onnx_session
from rag.schemas import RAGResponse, ScoredChunk

logger = logging.getLogger(__name__)


_SUPPORTED_EXTENSIONS = {
    ".md", ".txt", ".yaml", ".yml", ".json",
    ".pdf", ".docx", ".odt", ".odp", ".ods",
    ".csv", ".xlsx", ".xls", ".pptx",
    ".html", ".htm", ".zim",
}


def _read_file(path: Path) -> str:
    ext = path.suffix.lower()

    if ext == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)

    if ext == ".docx":
        from docx import Document as DocxDocument
        doc = DocxDocument(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)

    if ext in (".odt", ".odp", ".ods"):
        from odf import text, teletype
        from odf.opendocument import load
        doc = load(str(path))
        paragraphs = doc.getElementsByType(text.P)
        lines = [teletype.extractText(p) for p in paragraphs if teletype.extractText(p)]
        return "\n\n".join(lines)

    if ext == ".csv":
        import csv, io
        with open(path, newline="", encoding="utf-8", errors="replace") as f:
            rows = list(csv.reader(f))
        if not rows:
            return ""
        header = " | ".join(rows[0])
        data = "\n".join(" | ".join(row) for row in rows[1:])
        return f"{header}\n{data}"

    if ext in (".xlsx", ".xls"):
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cleaned = [str(c) if c is not None else "" for c in row]
                rows.append(" | ".join(cleaned))
            if rows:
                parts.append(f"[Sheet: {sheet}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)

    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(cells))
            if texts:
                slides.append(f"[Slide {slide_num}]\n" + "\n".join(texts))
        return "\n\n".join(slides)

    if ext in (".html", ".htm"):
        from html.parser import HTMLParser

        class _TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self._texts = []
                self._skip = False
            def handle_starttag(self, tag, attrs):
                if tag in ("script", "style"):
                    self._skip = True
            def handle_endtag(self, tag):
                if tag in ("script", "style"):
                    self._skip = False
            def handle_data(self, data):
                if not self._skip:
                    stripped = data.strip()
                    if stripped:
                        self._texts.append(stripped)
            def result(self):
                return "\n\n".join(self._texts)

        extractor = _TextExtractor()
        extractor.feed(path.read_text(encoding="utf-8", errors="replace"))
        return extractor.result()

    if ext == ".zim":
        import zim as zim_mod
        parts = []
        zim_file = zim_mod.open(str(path))
        try:
            for article in zim_file.articles:
                text = article.text if article.text else ""
                title = article.title if article.title else ""
                url = article.url if article.url else ""
                if text.strip():
                    header = f"[{title}]({url})" if title else ""
                    parts.append(f"{header}\n{text}" if header else text)
        finally:
            zim_file.close()
        return "\n\n---\n\n".join(parts) if parts else ""

    return path.read_text(encoding="utf-8", errors="replace")


class HakeemRAGEngine:
    def __init__(self, settings: HakeemRAGSettings):
        self._settings = settings
        self._embed_model: Optional[OnnxEmbedding] = None
        self._router: Optional[HakeemSemanticRouter] = None
        self._qdrant: Optional[HakeemQdrantStore] = None
        self._parent_retriever: Optional[HakeemParentRetriever] = None
        self._graph: Optional[HakeemKnowledgeGraph] = None
        self._decomposer: Optional[HakeemQueryDecomposer] = None
        self._hybrid_retriever: Optional[HakeemHybridRetriever] = None
        self._reranker: Optional[HakeemReranker] = None
        self._crag: Optional[HakeemCorrectiveRAG] = None
        self._citation: Optional[HakeemCitationFormatter] = None
        self._initialized = False
        self._init_lock = asyncio.Lock()
        self._mtime_cache: dict[str, float] = {}

    @property
    def is_initialized(self) -> bool:
        return self._initialized

    async def initialize(self):
        async with self._init_lock:
            if self._initialized:
                return
            await self._initialize_once()

    async def _initialize_once(self):
        s = self._settings
        loop = asyncio.get_event_loop()

        def _init_embedding():
            os.makedirs(s.model_dir, exist_ok=True)
            download_onnx_model(model_dir=s.model_dir)
            _load_onnx_session(s.model_dir, device=s.device)

            self._embed_model = OnnxEmbedding(
                model_dir=s.model_dir,
                device=s.device,
                embed_batch_size=s.indexing_batch_size,
            )

        await loop.run_in_executor(None, _init_embedding)
        logger.info("Embedding model loaded")

        self._qdrant = HakeemQdrantStore(
            host=s.qdrant_host,
            port=s.qdrant_port,
            api_key=s.qdrant_api_key,
            prefix=s.collection_name_prefix,
            vector_size=s.vector_size,
        )
        await self._qdrant.initialize(s.domains)
        logger.info("Qdrant store initialized")

        self._parent_retriever = HakeemParentRetriever(
            embed_model=self._embed_model,
            child_chunk_size=s.child_chunk_size,
            child_chunk_overlap=s.child_chunk_overlap,
            parent_chunk_size=s.parent_chunk_size,
            parent_chunk_overlap=s.parent_chunk_overlap,
        )
        logger.info("Parent retriever initialized")

        self._router = HakeemSemanticRouter(
            model_path=s.router_model_path,
            domains=s.domains,
            threshold=s.router_threshold,
            device=s.device,
        )
        await self._router.initialize()
        logger.info("Semantic router initialized")

        self._graph = HakeemKnowledgeGraph(
            uri=s.neo4j_uri,
            user=s.neo4j_user,
            password=s.neo4j_password,
            traversal_depth=s.graph_traversal_depth,
        )
        await self._graph.initialize()
        logger.info("Knowledge graph initialized")

        self._decomposer = HakeemQueryDecomposer(
            api_base=s.llm_api_base,
            model=s.llm_model,
            api_key=s.llm_api_key,
            num_queries=s.decomposer_num_queries,
        )
        logger.info("Query decomposer initialized")

        self._hybrid_retriever = HakeemHybridRetriever(
            qdrant_store=self._qdrant,
            rrf_k=s.rrf_k,
            top_k=s.hybrid_top_k,
        )

        self._reranker = HakeemReranker(
            model_path=s.reranker_model_path,
            device=s.reranker_device,
            top_k=s.reranker_top_k,
        )
        await self._reranker.initialize()
        logger.info("Reranker initialized")

        self._crag = HakeemCorrectiveRAG(
            api_base=s.llm_api_base,
            model=s.llm_model,
            api_key=s.llm_api_key,
            enabled=s.crag_enabled,
        )

        self._citation = HakeemCitationFormatter()

        self._initialized = True
        logger.info("HakeemRAGEngine fully initialized: %d domains, vector_size=%d",
                     len(s.domains), s.vector_size)

    async def query(self, user_message: str,
                    llm_api_base: Optional[str] = None,
                    llm_api_key: Optional[str] = None,
                    llm_model: Optional[str] = None) -> RAGResponse:
        if not self._initialized:
            return RAGResponse(error="RAG engine not initialized",
                               sufficient=False)

        s = self._settings

        api_base = llm_api_base or s.llm_api_base
        api_key = llm_api_key or s.llm_api_key
        model = llm_model or s.llm_model

        routes = await self._router.route(user_message)
        domains = [r.domain for r in routes]
        all_entities = []
        for r in routes:
            all_entities.extend(r.entities)

        logger.info("RAG route: domains=%s, entities=%d",
                     domains, len(all_entities))

        if s.decomposer_enabled and s.decomposer_num_queries > 1:
            try:
                sub_queries = await self._decomposer.decompose(
                    user_message, s.decomposer_num_queries,
                    api_base=api_base, api_key=api_key,
                )
            except Exception as e:
                logger.warning("Decomposition failed, using original query: %s", e)
                sub_queries = [user_message]
        else:
            sub_queries = [user_message]

        def _make_sparse(text: str) -> dict:
            tokens = text.lower().split()
            tf: dict[str, int] = {}
            for t in tokens:
                tf[t] = tf.get(t, 0) + 1
            max_freq = max(tf.values()) if tf else 1
            sparse_indices = [abs(hash(t)) % 100000 for t in tf]
            sparse_values = [freq / max_freq for freq in tf.values()]
            return {"indices": sparse_indices, "values": sparse_values}

        query_embeddings: list[tuple[str, list[float], dict]] = []
        for sq in sub_queries:
            dv = await self._embed_model._aget_query_embedding(sq)
            sv = _make_sparse(sq)
            query_embeddings.append((sq, dv, sv))

        scored_chunks = await self._hybrid_retriever.search(
            query=user_message,
            domains=domains,
            dense_vector=query_embeddings[0][1],
            sparse_vector=query_embeddings[0][2],
            query_embeddings=query_embeddings,
        )

        if not scored_chunks:
            logger.info("No vector results for query")
            return RAGResponse(sufficient=False,
                               verification_status="no_results",
                               domains=domains)

        parent_chunks = self._parent_retriever.resolve_parents(scored_chunks)

        top_chunks = await self._reranker.rerank(user_message, parent_chunks)

        graph_paths = await self._graph.graph_traversal(all_entities)

        crag_result = await self._crag.verify(user_message, top_chunks)

        if not crag_result.sufficient and crag_result.status == "insufficient":
            logger.warning("CRAG: insufficient context for query")
            abstention = await self._crag.abstain()
            return RAGResponse(
                context=abstention,
                formatted_context=abstention,
                chunks=top_chunks,
                graph_paths=graph_paths,
                domains=domains,
                verification_status="insufficient",
                sufficient=False,
            )

        rag_response = self._citation.build_rag_response(
            chunks=top_chunks,
            graph_paths=graph_paths,
            verification_status=crag_result.status,
            sufficient=crag_result.sufficient,
        )

        return rag_response

    async def index_document(self, file_path: str, domain: str,
                              doc_id: Optional[str] = None) -> dict[str, Any]:
        if not self._initialized:
            raise RuntimeError("HakeemRAGEngine not initialized")

        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        if domain not in self._settings.domains:
            raise ValueError(f"Unknown domain '{domain}'. Valid: {self._settings.domains}")

        doc_id = doc_id or path.stem
        text = _read_file(path)

        if not text.strip():
            logger.warning("Empty file: %s", file_path)
            return {"doc_id": doc_id, "filename": path.name,
                    "chunks": 0, "status": "empty", "domain": domain}

        mtime = path.stat().st_mtime
        filename = path.name
        source_file = str(path.resolve())

        child_chunks, child_embeddings = self._parent_retriever.chunk_document(
            text=text, doc_id=doc_id, filename=filename,
            source_file=source_file, domain=domain, mtime=mtime,
        )

        if not child_chunks:
            return {"doc_id": doc_id, "filename": path.name,
                    "chunks": 0, "status": "empty", "domain": domain}

        await self._qdrant.delete_document(domain, doc_id)
        await self._qdrant.add_chunks(domain, child_chunks, child_embeddings)

        entities = await self._router.extract_entities(text)
        await self._graph.index_document(doc_id, text, domain, entities)

        self._mtime_cache[source_file] = mtime

        logger.info("Indexed %s into domain=%s: %d chunks, %d entities",
                     filename, domain, len(child_chunks), len(entities))
        return {
            "doc_id": doc_id,
            "filename": filename,
            "chunks": len(child_chunks),
            "status": "indexed",
            "domain": domain,
        }

    async def index_directory(self, directory: str,
                               domain: str) -> list[dict[str, Any]]:
        if not self._initialized:
            raise RuntimeError("HakeemRAGEngine not initialized")

        path = Path(directory)
        if not path.exists():
            logger.warning("Directory not found: %s", directory)
            return []

        sem = asyncio.Semaphore(4)
        results = []

        async def _index_one(f: Path):
            async with sem:
                return await self.index_document(str(f), domain)

        for f in sorted(path.rglob("*")):
            if f.is_file() and f.suffix.lower() in _SUPPORTED_EXTENSIONS:
                try:
                    result = await _index_one(f)
                    results.append(result)
                    delay = self._settings.indexing_delay_ms / 1000
                    if delay > 0:
                        await asyncio.sleep(delay)
                except Exception as e:
                    logger.error("Failed to index %s: %s", f, e)
        logger.info("Indexed %d files from %s (domain=%s)",
                     len(results), directory, domain)
        return results

    async def index_if_changed(self, dirs: dict[str, str]) -> int:
        if not self._initialized:
            return 0

        total_chunks = 0
        for domain, directory in dirs.items():
            if domain not in self._settings.domains:
                logger.warning("Unknown domain '%s', skipping", domain)
                continue

            dpath = Path(directory)
            if not dpath.exists():
                logger.warning("Source directory not found: %s", directory)
                continue

            for f in sorted(dpath.rglob("*")):
                if not f.is_file() or f.suffix.lower() not in _SUPPORTED_EXTENSIONS:
                    continue

                current_mtime = f.stat().st_mtime
                cached = self._mtime_cache.get(str(f))
                if cached is not None and current_mtime <= cached:
                    continue

                result = await self.index_document(str(f), domain)
                total_chunks += result.get("chunks", 0)
                self._mtime_cache[str(f)] = current_mtime

        logger.info("index_if_changed done: %d chunks indexed", total_chunks)
        return total_chunks

    async def search(self, query: str,
                     domain: Optional[str] = None) -> RAGResponse:
        return await self.query(query)

    async def delete_document(self, domain: str, doc_id: str) -> bool:
        return await self._qdrant.delete_document(domain, doc_id)

    async def list_documents(self) -> list[dict[str, Any]]:
        all_docs: list[dict] = []
        for domain in self._settings.domains:
            docs = await self._qdrant.list_documents(domain)
            all_docs.extend(docs)
        return all_docs

    async def close(self):
        if self._graph:
            await self._graph.close()


_engine: Optional[HakeemRAGEngine] = None


async def get_rag_engine() -> Optional[HakeemRAGEngine]:
    return _engine


async def init_rag_engine(settings: HakeemRAGSettings) -> Optional[HakeemRAGEngine]:
    global _engine
    if not settings.enabled:
        logger.info("HakeemRAG is disabled")
        return None
    engine = HakeemRAGEngine(settings)
    await engine.initialize()
    _engine = engine
    return engine
